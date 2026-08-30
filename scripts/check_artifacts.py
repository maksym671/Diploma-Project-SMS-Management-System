"""Check the deck and the documentation against the repository they describe.

The slides quote counts — tests, templates, lines, translatable strings — that
drift the moment the project moves on. This recomputes each one and fails if a
built artefact still claims the old number, so "the presentation is finished"
survives the next commit.

    ./pptx_env/bin/python scripts/check_artifacts.py
"""
import pathlib
import re
import sys

ROOT = pathlib.Path(__file__).resolve().parent.parent


def facts():
    """Every number the artefacts are allowed to quote, measured from source."""
    po = (ROOT / 'locale/pl/LC_MESSAGES/django.po').read_text(encoding='utf-8')
    return {
        'tests': (ROOT / 'core/tests.py').read_text().count('\n    def test_'),
        'templates': len(list((ROOT / 'templates').rglob('*.html'))),
        'css lines': len((ROOT / 'static/css/style.css').read_text().splitlines()),
        'js lines': len((ROOT / 'static/js/main.js').read_text().splitlines()),
        'migrations': len([
            p for p in (ROOT / 'core/migrations').glob('*.py')
            if p.name != '__init__.py'
        ]),
        'po entries': len(re.findall(
            r'\nmsgid ((?:"(?:[^"\\]|\\.)*"\n?)+)msgstr', '\n' + po,
        )) - 1,
        'fuzzy': po.count('#, fuzzy'),
    }


def artefact_text():
    """The text of everything that leaves the repository."""
    from pptx import Presentation
    import docx

    texts = {}
    deck = ROOT / 'Maksym_Shpak_Diploma_Presentation.pptx'
    if deck.exists():
        texts['deck.pptx'] = '\n'.join(
            shape.text_frame.text
            for slide in Presentation(deck).slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
    doc = ROOT / 'docs/SMS_Diploma_Documentation_Shpak_Maksym.docx'
    if doc.exists():
        texts['doc.docx'] = '\n'.join(p.text for p in docx.Document(doc).paragraphs)
    for name in ('docs/SMS_Diploma_Documentation_Shpak_Maksym.md', 'README.md'):
        path = ROOT / name
        if path.exists():
            texts[name] = path.read_text(encoding='utf-8')
    return texts


def main():
    measured = facts()
    print('Measured from the repository:')
    for key, value in measured.items():
        print(f'  {key:12} {value}')

    # A claim is stale when an artefact quotes a count near the real one but
    # not equal to it — "96 tests" once the suite runs 111.
    patterns = {
        'tests': (r'(\d+)\s+(?:automated\s+)?tests', measured['tests']),
        'templates': (r'Django Templates · (\d+) files', measured['templates']),
        'js lines': (r'JavaScript ES6 · (\d+) lines', measured['js lines']),
        'css lines': (r'custom properties · ([\d\s]+) lines', measured['css lines']),
        'po entries': (r'(\d+)\s+(?:translatable\s+)?(?:strings|messages)',
                       measured['po entries']),
    }

    problems = []
    for name, text in artefact_text().items():
        for label, (pattern, expected) in patterns.items():
            for found in re.findall(pattern, text):
                if int(found.replace(' ', '').replace(' ', '')) != expected:
                    problems.append(f'{name}: claims {found} {label}, repository has {expected}')

    if measured['fuzzy']:
        problems.append(
            f"locale/pl: {measured['fuzzy']} fuzzy entries — msgfmt drops them "
            'and they ship as English'
        )

    print()
    if problems:
        print('STALE CLAIMS:')
        for problem in problems:
            print(f'  {problem}')
        return 1
    print('Every quoted number matches the repository.')
    return 0


if __name__ == '__main__':
    sys.exit(main())
