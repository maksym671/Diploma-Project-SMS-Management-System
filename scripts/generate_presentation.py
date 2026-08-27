"""Builds the diploma defence deck for the Student Management System.

Every figure quoted on the slides comes from the repository itself: 6 models,
6 migrations, 80 tests, 308 translatable strings, and the demo dataset that
ships with the deployment.

    pptx_env/bin/python scripts/generate_presentation.py
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / 'docs' / 'screenshots'
OUTPUT = ROOT / 'Maksym_Shpak_Diploma_Presentation.pptx'

LIVE_URL = 'diploma-project-sms-management-system.onrender.com'
LIVE_HREF = 'https://' + LIVE_URL
UNI_URL = 'https://www.vizja.pl'
DEMO_TEACHER = 'prof.martinez'
DEMO_TEACHER_B = 'prof.chen'
DEMO_PASSWORD = 'demo1234'
STUDENT = 'Maksym Shpak'
INDEX_NO = '45567'
SUPERVISOR = 'Marcin Kacprowicz'
LOGO_DARK = ROOT / 'docs' / 'assets' / 'vizja_logo_dark.png'
LOGO_LIGHT = ROOT / 'docs' / 'assets' / 'vizja_logo_light.png'
LOGO_DARK_V = ROOT / 'docs' / 'assets' / 'vizja_logo_dark_vertical.png'

# --- palette -----------------------------------------------------------------
INK = RGBColor(0x0F, 0x17, 0x2A)
INK_SOFT = RGBColor(0x1E, 0x29, 0x3B)
SLATE = RGBColor(0x47, 0x55, 0x69)
MUTED = RGBColor(0x8C, 0x9A, 0xAF)
FAINT = RGBColor(0xB4, 0xC0, 0xD0)
LINE = RGBColor(0xE3, 0xE9, 0xF0)
CANVAS = RGBColor(0xF6, 0xF8, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

ACCENT = RGBColor(0x25, 0x63, 0xEB)
ACCENT_DEEP = RGBColor(0x1D, 0x4E, 0xD8)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
PINK = RGBColor(0xDB, 0x27, 0x77)
TEAL = RGBColor(0x0D, 0x94, 0x88)

BODY_FONT = 'Calibri'
HEAD_FONT = 'Calibri'

W, H = 13.3333, 7.5
MARGIN = 0.62
CONTENT_W = W - 2 * MARGIN
BODY_TOP = 1.62
FOOTER_Y = 6.94

TOTAL_SLIDES = 11


# --- primitives --------------------------------------------------------------
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def paint(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill=None, line=None, radius=0.06,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = radius
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return sh


def card(slide, x, y, w, h, fill=WHITE, line=LINE, radius=0.05):
    return rect(slide, x, y, w, h, fill=fill, line=line, radius=radius)


def track(slide, x, y, w, h, color=ACCENT, radius=0.5):
    """A slim rounded accent bar used as a visual marker."""
    return rect(slide, x, y, w, h, fill=color, radius=radius)


def _style_run(run, size=12, bold=False, color=SLATE, font=BODY_FONT,
               italic=False, spacing=None, url=None):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    if spacing:
        run._r.get_or_add_rPr().set('spc', str(int(spacing * 100)))
    if url:
        run.hyperlink.address = url
        run.font.underline = False
        run.font.color.rgb = color


def text(slide, x, y, w, h, blocks, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """Draws a text box. Each block is a dict describing one paragraph.

    A block carries either ``text`` (single run) or ``runs`` (list of
    ``(string, style)`` pairs) plus optional paragraph settings.
    """
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, block in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = block.get('align', align)
        if 'space_before' in block:
            p.space_before = Pt(block['space_before'])
        p.space_after = Pt(block.get('space_after', 0))
        if 'line' in block:
            p.line_spacing = block['line']
        runs = block.get('runs')
        if runs is None:
            runs = [(block.get('text', ''), block)]
        for content, style in runs:
            run = p.add_run()
            run.text = content
            _style_run(
                run,
                size=style.get('size', 12),
                bold=style.get('bold', False),
                color=style.get('color', SLATE),
                font=style.get('font', BODY_FONT),
                italic=style.get('italic', False),
                spacing=style.get('spacing'),
                url=style.get('url'),
            )
    return box


# --- slide chrome ------------------------------------------------------------
def add_logo(slide, path, x, y, w, url=UNI_URL):
    """Places a PNG logo and makes it open the university site on click."""
    iw, ih = Image.open(path).size
    h = w * ih / iw
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    pic.click_action.hyperlink.address = url
    return h


def footer(slide, dark=False, logo=False):
    if logo:
        path = LOGO_DARK if dark else LOGO_LIGHT
        add_logo(slide, path, MARGIN, FOOTER_Y - 0.02, 1.42)
    text(
        slide, W - MARGIN - 2.4, FOOTER_Y, 2.4, 0.28,
        [{'text': 'www.vizja.pl', 'size': 10.5,
          'color': FAINT if dark else MUTED, 'spacing': 0.6,
          'url': UNI_URL}],
        align=PP_ALIGN.RIGHT,
    )


def head(prs, number, title, kicker=None, accent=ACCENT):
    """Standard light content slide with title, section marker and counter."""
    slide = blank(prs)
    paint(slide, CANVAS)
    track(slide, MARGIN, 0.62, 0.11, 0.46, color=accent)
    y = 0.55
    if kicker:
        text(
            slide, MARGIN + 0.3, y, 7.0, 0.24,
            [{'text': kicker.upper(), 'size': 10.5, 'bold': True,
              'color': accent, 'spacing': 1.4}],
        )
        y += 0.28
    text(
        slide, MARGIN + 0.3, y, 9.6, 0.6,
        [{'text': title, 'size': 30, 'bold': True, 'color': INK, 'font': HEAD_FONT}],
    )
    text(
        slide, W - MARGIN - 2.0, 0.6, 2.0, 0.3,
        [{'text': f'{number} / {TOTAL_SLIDES}', 'size': 11.5,
          'color': MUTED, 'spacing': 0.8}],
        align=PP_ALIGN.RIGHT,
    )
    rect(slide, MARGIN, FOOTER_Y - 0.16, CONTENT_W, 0.012,
         fill=LINE, shape=MSO_SHAPE.RECTANGLE)
    footer(slide, logo=True)
    return slide


def eyebrow(slide, x, y, w, label, color=MUTED):
    text(slide, x, y, w, 0.22,
         [{'text': label.upper(), 'size': 10, 'bold': True,
           'color': color, 'spacing': 1.2}])


def bullets(slide, x, y, w, items, size=12.5, gap=7.2, color=SLATE,
            lead_color=INK, marker=None, marker_color=ACCENT):
    """Bullet lines. An item may be a string or a (lead, rest) tuple."""
    blocks = []
    for item in items:
        if isinstance(item, tuple):
            lead, rest = item
        else:
            lead, rest = None, item
        runs = []
        if marker:
            runs.append((marker + '  ', {'size': size, 'bold': True,
                                         'color': marker_color}))
        if lead:
            runs.append((lead + ' ', {'size': size, 'bold': True, 'color': lead_color}))
        runs.append((rest, {'size': size, 'color': color}))
        blocks.append({'runs': runs, 'space_after': gap, 'line': 1.16})
    return text(slide, x, y, w, 4.6, blocks)


def stat_tile(slide, x, y, w, h, value, label, color=ACCENT, fill=WHITE,
              value_size=30, label_size=10.5):
    card(slide, x, y, w, h, fill=fill)
    track(slide, x, y + 0.22, 0.055, h - 0.44, color=color, radius=0.5)
    text(slide, x + 0.3, y + 0.2, w - 0.5, h - 0.4, [
        {'text': value, 'size': value_size, 'bold': True, 'color': color,
         'font': HEAD_FONT, 'space_after': 2},
        {'text': label, 'size': label_size, 'color': SLATE, 'line': 1.1},
    ], anchor=MSO_ANCHOR.MIDDLE)


def panel(slide, x, y, w, h, title, items, accent=ACCENT, size=12.2, gap=7.0,
          note=None):
    card(slide, x, y, w, h)
    rect(slide, x, y, w, 0.055, fill=accent, radius=0.5)
    eyebrow(slide, x + 0.28, y + 0.3, w - 0.5, title, color=accent)
    bullets(slide, x + 0.28, y + 0.66, w - 0.56, items, size=size, gap=gap,
            marker='—', marker_color=FAINT)
    if note:
        text(slide, x + 0.28, y + h - 0.5, w - 0.56, 0.34,
             [{'text': note, 'size': 10.5, 'italic': True, 'color': MUTED,
               'line': 1.1}])


def note_bar(slide, y, runs, fill=INK, height=0.62):
    bar = rect(slide, MARGIN, y, CONTENT_W, height, fill=fill, radius=0.09)
    text(slide, MARGIN + 0.34, y + 0.06, CONTENT_W - 0.68, height - 0.12,
         [{'runs': runs, 'line': 1.12}], anchor=MSO_ANCHOR.MIDDLE)
    return bar


def clean_table(gframe, header_fill=INK, header_color=WHITE,
                zebra=RGBColor(0xF2, 0xF5, 0xF9)):
    table = gframe.table
    tbl = table._tbl
    tblPr = tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')
        for child in list(tblPr):
            if child.tag == qn('a:tableStyleId'):
                tblPr.remove(child)
    for r, row in enumerate(table.rows):
        for cell in row.cells:
            cell.margin_left = Inches(0.16)
            cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.07)
            cell.margin_bottom = Inches(0.07)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = header_fill
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else zebra
    return table


def table_text(cell, value, size=10.5, bold=False, color=SLATE, spacing=None,
               align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.06
    run = p.add_run()
    run.text = value
    _style_run(run, size=size, bold=bold, color=color, spacing=spacing)


def matrix(slide, x, y, w, headers, rows, col_widths, row_height=0.395,
           header_height=0.42, highlight=None, highlight_color=ACCENT):
    gframe = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w),
        Inches(header_height + row_height * len(rows)),
    )
    table = clean_table(gframe)
    table.rows[0].height = Inches(header_height)
    for i, cw in enumerate(col_widths):
        table.columns[i].width = Inches(cw)
    for i, label in enumerate(headers):
        table_text(table.cell(0, i), label.upper(), size=9.5, bold=True,
                   color=WHITE, spacing=0.9)
    for r, row in enumerate(rows, start=1):
        table.rows[r].height = Inches(row_height)
        for c, value in enumerate(row):
            is_key = c == 0
            is_hi = highlight is not None and c == highlight
            table_text(
                table.cell(r, c), value,
                size=10.2,
                bold=is_key or is_hi,
                color=INK if is_key else (highlight_color if is_hi else SLATE),
            )
    return table


def picture_frame(slide, x, y, w, path, caption=None, aspect=1.6):
    h = w / aspect
    card(slide, x - 0.055, y - 0.055, w + 0.11, h + 0.11, fill=WHITE, line=LINE)
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    if caption:
        text(slide, x, y + h + 0.14, w, 0.24,
             [{'text': caption, 'size': 10, 'color': MUTED, 'spacing': 0.5}])
    return h


# --- slides ------------------------------------------------------------------
def slide_title(prs):
    slide = blank(prs)
    paint(slide, INK)
    rect(slide, 8.55, -1.1, 6.6, 6.6, fill=INK_SOFT, radius=0.5)
    rect(slide, 10.35, 3.5, 4.4, 4.4, fill=RGBColor(0x18, 0x22, 0x39), radius=0.5)
    track(slide, MARGIN, 1.42, 0.11, 0.9, color=ACCENT)

    add_logo(slide, LOGO_DARK_V, 9.25, 1.72, 3.55)

    text(slide, MARGIN + 0.32, 1.44, 8.0, 0.26,
         [{'text': 'UNIWERSYTET VIZJA  ·  SCHOOL OF COMPUTER SCIENCE & TECHNOLOGIES',
           'size': 11, 'bold': True, 'color': FAINT, 'spacing': 1.4}])
    text(slide, MARGIN + 0.32, 1.86, 8.2, 2.0, [
        {'text': 'Design and Implementation of a', 'size': 33,
         'color': WHITE, 'font': HEAD_FONT, 'line': 1.1},
        {'text': 'Web-Based Student Management System', 'size': 33, 'bold': True,
         'color': WHITE, 'font': HEAD_FONT, 'line': 1.1},
        {'text': 'for Educational Institutions', 'size': 33,
         'color': WHITE, 'font': HEAD_FONT, 'line': 1.1},
    ])

    chip = rect(slide, MARGIN + 0.32, 4.14, 6.35, 0.44, fill=ACCENT, radius=0.4)
    chip.click_action.hyperlink.address = LIVE_HREF
    text(slide, MARGIN + 0.52, 4.2, 6.05, 0.32,
         [{'runs': [('LIVE   ', {'size': 9.5, 'bold': True, 'color': WHITE,
                                 'spacing': 1.4}),
                    (LIVE_HREF, {'size': 11.5, 'bold': True, 'color': WHITE,
                                 'url': LIVE_HREF})]}],
         anchor=MSO_ANCHOR.MIDDLE)

    rect(slide, MARGIN, 5.42, 7.9, 0.012, fill=RGBColor(0x2C, 0x3A, 0x52),
         shape=MSO_SHAPE.RECTANGLE)
    meta = [
        ('STUDENT', STUDENT, f'Index: {INDEX_NO}'),
        ('SUPERVISOR', SUPERVISOR, 'Diploma seminar'),
        ('FIELD OF STUDY', 'Computer Science', 'Engineering degree · 2026'),
    ]
    for i, (label, value, sub) in enumerate(meta):
        x = MARGIN + i * 2.66
        text(slide, x, 5.66, 2.5, 0.9, [
            {'text': label, 'size': 9.5, 'bold': True, 'color': ACCENT,
             'spacing': 1.2, 'space_after': 5},
            {'text': value, 'size': 14, 'bold': True, 'color': WHITE,
             'space_after': 2},
            {'text': sub, 'size': 10.5, 'color': MUTED},
        ])
    footer(slide, dark=True)
    return slide


def slide_agenda(prs):
    slide = head(prs, 2, 'Presentation Agenda', kicker='What this defence covers')
    items = [
        ('01', 'Problem Statement & Purpose',
         'Spreadsheets, paper registers, no audit trail'),
        ('02', 'Comparative Analysis',
         'Moodle, USOS, Google Classroom — and the gap'),
        ('03', 'Technology Stack',
         'Django 6, Python 3.12, PostgreSQL 16'),
        ('04', 'System Architecture',
         'Browser › Render TLS › Gunicorn › PostgreSQL'),
        ('05', 'Database Design',
         '6 models, 6 migrations, constraints in the schema'),
        ('06', 'Key Features',
         'Weighted grades, RBAC, EN/PL, 80 tests'),
        ('07', 'Application Screenshots',
         'Dashboard, student records, mobile layout'),
        ('08', 'Deployment & Results',
         'Live on Render + Neon, then what comes next'),
    ]
    col_w = (CONTENT_W - 0.4) / 2
    for i, (num, title, sub) in enumerate(items):
        col, row = divmod(i, 4)
        x = MARGIN + col * (col_w + 0.4)
        y = BODY_TOP + 0.08 + row * 1.22
        text(slide, x, y + 0.04, 0.52, 0.4,
             [{'text': num, 'size': 22, 'bold': True, 'color': FAINT,
               'font': HEAD_FONT}])
        text(slide, x + 0.62, y, col_w - 0.62, 0.95, [
            {'text': title, 'size': 16, 'bold': True, 'color': INK,
             'space_after': 4},
            {'text': sub, 'size': 12, 'color': SLATE, 'line': 1.14},
        ])
        if row < 3:
            rect(slide, x, y + 1.04, col_w, 0.008, fill=LINE,
                 shape=MSO_SHAPE.RECTANGLE)
    return slide


def slide_problem(prs):
    slide = head(prs, 3, 'Problem Statement & Purpose', kicker='Why this system')
    tiles = [
        ('4–6 h', 'per week per teacher spent on\nmanual record-keeping', ACCENT),
        ('3 places', 'a single student record is\nsplit across today', AMBER),
        ('0', 'audit trail of who entered\nor edited a grade', PINK),
    ]
    tw = (CONTENT_W - 0.44) / 3
    for i, (value, label, color) in enumerate(tiles):
        stat_tile(slide, MARGIN + i * (tw + 0.22), BODY_TOP, tw, 1.12,
                  value, label, color=color, value_size=26, label_size=10.5)

    panel_w = (CONTENT_W - 0.28) / 2
    panel(
        slide, MARGIN, BODY_TOP + 1.34, panel_w, 3.28,
        'The problem', [
            ('Fragmented records —', 'spreadsheets, paper registers, email threads'),
            ('No single view —', 'grades and attendance never meet in one place'),
            ('Manual arithmetic —', 'averages recalculated by hand, errors slip through'),
            ('No accountability —', 'nobody can prove who recorded a mark'),
            ('No confidentiality —', 'every teacher can see every student'),
            ('No Polish interface —', 'staff work in English-only tools'),
        ], accent=PINK, size=12.2, gap=7.6,
    )
    panel(
        slide, MARGIN + panel_w + 0.28, BODY_TOP + 1.34, panel_w, 3.28,
        'The purpose', [
            ('One source of truth —', 'a relational schema with enforced constraints'),
            ('Instant analytics —', 'KPI tiles and three charts from ORM aggregation'),
            ('Automatic authorship —', 'every grade stores who assigned it'),
            ('Scoped visibility —', 'teachers see their own courses, admins see all'),
            ('Bilingual by design —', 'the whole interface switches EN ⇄ PL'),
            ('Deployable —', 'one build script, HTTPS, managed PostgreSQL'),
        ], accent=GREEN, size=12.2, gap=7.6,
    )
    note_bar(slide, FOOTER_Y - 0.98, [
        ('Result: ', {'size': 12.5, 'bold': True, 'color': WHITE}),
        ('one auditable, role-aware platform now managing 25 students, 8 courses, '
         '87 enrollments and 87 grades in production.',
         {'size': 12.5, 'color': RGBColor(0xD6, 0xDF, 0xEC)}),
    ])
    return slide


def slide_competitors(prs):
    slide = head(prs, 4, 'Comparative Analysis', kicker='Existing platforms')
    headers = ['Criterion', 'Moodle', 'USOS', 'Google Classroom',
               'SMS — this project']
    rows = [
        ['Target user', 'Large universities', 'Polish public HEIs',
         'Schools and courses', 'Departments and small institutions'],
        ['Deployment', 'Self-hosted, heavy', 'Vendor-managed',
         'Google cloud only', 'Self-hosted, one Django app'],
        ['Setup effort', 'Days of plugin work', 'Institutional contract',
         'Minutes, no control', 'Minutes — a single build script'],
        ['Gradebook', 'Yes, complex', 'Yes', 'Basic',
         'Weighted components + author audit'],
        ['Attendance register', 'Plugin required', 'Yes', 'Not available',
         'Built in, three statuses'],
        ['Role isolation', 'Configurable, intricate', 'Yes', 'Limited',
         'Enforced at queryset level'],
        ['Bilingual EN / PL', 'Language packs', 'Polish-first', 'Yes',
         'Native, 308 strings'],
        ['Analytics', 'Reports plugin', 'Institutional reports', 'Minimal',
         'Chart.js dashboard + CSV'],
        ['Cost of ownership', 'GPL + hosting', 'Licensed', 'Free, data in Google',
         'Open source, €0 hosting'],
    ]
    widths = [2.28, 2.16, 2.16, 2.36, 3.13]
    matrix(slide, MARGIN, BODY_TOP - 0.06, CONTENT_W, headers, rows, widths,
           row_height=0.435, highlight=4)
    text(slide, MARGIN, FOOTER_Y - 0.52, CONTENT_W, 0.3,
         [{'runs': [
             ('Positioning: ', {'size': 11.5, 'bold': True, 'color': INK}),
             ('the established platforms are either too heavy or too closed for a '
              'single department; SMS trades breadth for a focused, auditable and '
              'fully self-hosted core.',
              {'size': 11.5, 'color': SLATE})]}])
    return slide


def slide_stack(prs):
    slide = head(prs, 5, 'Technology Stack', kicker='Chosen for a maintainable monolith')
    groups = [
        ('Backend', 'Runtime · data · authentication', ACCENT, [
            'Python 3.12', 'Django 6.0.6', 'PostgreSQL 16 (Neon)',
            'SQLite 3 for development', 'Django session authentication',
            'Custom user model with roles', 'Django ORM · 6 migrations',
            'dj-database-url 2.1',
        ]),
        ('Frontend', 'Templates · charts · theming', VIOLET, [
            'Django Templates · 27 files', 'Semantic HTML5',
            'CSS custom properties · 1 945 lines', 'Vanilla JavaScript ES6 · 413 lines',
            'Chart.js 4 — doughnut and bars', 'Turbo 8 page transitions',
            'Bootstrap Icons', 'Light and dark themes',
        ]),
        ('Infrastructure', 'Hosting · operations · CI', GREEN, [
            'Render web service', 'Neon managed PostgreSQL',
            'Gunicorn 22 (WSGI)', 'WhiteNoise 6.6 static files',
            'HTTPS with HSTS preload', 'GitHub Actions CI',
            'Git version control', 'Zero CDN dependencies',
        ]),
    ]
    cw = (CONTENT_W - 0.44) / 3
    for i, (title, sub, color, items) in enumerate(groups):
        x = MARGIN + i * (cw + 0.22)
        card(slide, x, BODY_TOP, cw, 4.42)
        rect(slide, x, BODY_TOP, cw, 0.055, fill=color, radius=0.5)
        text(slide, x + 0.3, BODY_TOP + 0.32, cw - 0.6, 0.66, [
            {'text': title.upper(), 'size': 12, 'bold': True, 'color': color,
             'spacing': 1.3, 'space_after': 4},
            {'text': sub, 'size': 10.5, 'color': MUTED},
        ])
        blocks = []
        for item in items:
            blocks.append({'runs': [
                ('·  ', {'size': 12.5, 'bold': True, 'color': color}),
                (item, {'size': 12.5, 'color': INK_SOFT}),
            ], 'space_after': 9.5, 'line': 1.1})
        text(slide, x + 0.3, BODY_TOP + 1.16, cw - 0.6, 3.0, blocks)
    note_bar(slide, FOOTER_Y - 0.86, [
        ('No build step, no separate API server. ',
         {'size': 12, 'bold': True, 'color': WHITE}),
        ('Every asset is vendored locally, so the application runs and demonstrates '
         'fully offline.', {'size': 12, 'color': RGBColor(0xD6, 0xDF, 0xEC)}),
    ], height=0.58)
    return slide


def slide_architecture(prs):
    slide = head(prs, 6, 'System Architecture', kicker='Server-rendered MTV monolith')
    stages = [
        ('CLIENT', 'Browser', 'HTML · CSS · ES6\nTurbo navigation', ACCENT),
        ('EDGE', 'Render · TLS 443', 'HTTPS, HSTS\nstatic via WhiteNoise', TEAL),
        ('APPLICATION', 'Gunicorn + Django 6', 'URL › View › Template\nORM aggregation', VIOLET),
        ('DATA', 'PostgreSQL 16', 'Neon managed\ndaily backups', GREEN),
    ]
    bw = 2.82
    gap = (CONTENT_W - 4 * bw) / 3
    for i, (label, title, sub, color) in enumerate(stages):
        x = MARGIN + i * (bw + gap)
        card(slide, x, BODY_TOP + 0.12, bw, 1.62)
        rect(slide, x, BODY_TOP + 0.12, bw, 0.055, fill=color, radius=0.5)
        text(slide, x + 0.26, BODY_TOP + 0.42, bw - 0.52, 1.1, [
            {'text': label, 'size': 9.5, 'bold': True, 'color': color,
             'spacing': 1.3, 'space_after': 5},
            {'text': title, 'size': 15, 'bold': True, 'color': INK,
             'font': HEAD_FONT, 'space_after': 5},
            {'text': sub, 'size': 10.5, 'color': SLATE, 'line': 1.14},
        ])
        if i < 3:
            text(slide, x + bw + 0.01, BODY_TOP + 0.72, gap - 0.02, 0.4,
                 [{'text': '›', 'size': 22, 'bold': True, 'color': FAINT}],
                 align=PP_ALIGN.CENTER)

    mid_y = BODY_TOP + 2.02
    card(slide, MARGIN, mid_y, CONTENT_W, 0.92, fill=WHITE)
    eyebrow(slide, MARGIN + 0.3, mid_y + 0.19, 4.0, 'Request pipeline — middleware')
    chips = ['Security', 'Session', 'CSRF', 'Locale', 'WhiteNoise',
             'Authentication', 'Role decorators', 'Queryset scoping']
    cx = MARGIN + 0.3
    for name in chips:
        cw = 0.2 + len(name) * 0.083
        rect(slide, cx, mid_y + 0.49, cw, 0.3, fill=CANVAS, line=LINE, radius=0.4)
        text(slide, cx, mid_y + 0.52, cw, 0.24,
             [{'text': name, 'size': 10, 'bold': True, 'color': INK_SOFT}],
             align=PP_ALIGN.CENTER)
        cx += cw + 0.12

    sup_y = mid_y + 1.12
    services = [
        ('i18n EN / PL', 'Interface language switch'),
        ('GitHub Actions', 'Tests and deploy gates'),
        ('Teachers screen', 'Admin issues lecturer accounts'),
        ('CSV export', 'Role-scoped reporting'),
        ('/api/dashboard/', 'JSON metrics endpoint'),
    ]
    sw = (CONTENT_W - 4 * 0.18) / 5
    eyebrow(slide, MARGIN, sup_y, 5.0, 'Supporting services')
    for i, (name, sub) in enumerate(services):
        x = MARGIN + i * (sw + 0.18)
        card(slide, x, sup_y + 0.3, sw, 0.86, fill=WHITE)
        text(slide, x + 0.22, sup_y + 0.46, sw - 0.44, 0.6, [
            {'text': name, 'size': 12, 'bold': True, 'color': INK,
             'space_after': 3},
            {'text': sub, 'size': 10, 'color': MUTED, 'line': 1.1},
        ])
    return slide


def slide_database(prs):
    slide = head(prs, 7, 'Database Design', kicker='6 models · 6 migrations')
    entities = [
        ('User', 'extends AbstractUser', ACCENT,
         ['id', 'username · email', 'first_name · last_name',
          'role → admin | teacher', 'password (PBKDF2)', 'is_active']),
        ('Student', 'academic identity', VIOLET,
         ['id', 'student_number (unique)', 'first_name · last_name', 'email',
          'study_program', 'date_enrolled · is_active']),
        ('Course', 'taught unit', TEAL,
         ['id', 'course_code (unique)', 'course_name · description',
          'semester · credits', 'max_students', 'teacher_id → User']),
        ('Enrollment', 'student ⇄ course', AMBER,
         ['id', 'student_id → Student', 'course_id → Course', 'enrollment_date',
          'status → active | completed | dropped', 'unique (student, course)']),
        ('Grade', 'weighted component', GREEN,
         ['id', 'enrollment_id → Enrollment',
          'kind → coursework | midterm | final | retake',
          'weight 1–100 %', 'grade_value 2.0 – 5.0', 'assigned_by → User']),
        ('Attendance', 'daily register', PINK,
         ['id', 'enrollment_id → Enrollment', 'date',
          'status → present | absent | late', 'notes',
          'unique (enrollment, date)']),
    ]
    cw = (CONTENT_W - 2 * 0.22) / 3
    ch = 2.16
    for i, (name, sub, color, fields) in enumerate(entities):
        row, col = divmod(i, 3)
        x = MARGIN + col * (cw + 0.22)
        y = BODY_TOP + row * (ch + 0.22)
        card(slide, x, y, cw, ch)
        rect(slide, x, y, 0.055, ch, fill=color, radius=0.5)
        text(slide, x + 0.28, y + 0.2, cw - 0.5, 0.5, [
            {'runs': [(name, {'size': 15, 'bold': True, 'color': INK,
                              'font': HEAD_FONT}),
                      ('   ' + sub, {'size': 10, 'italic': True, 'color': MUTED})]},
        ])
        blocks = [{'text': f_, 'size': 10.8, 'color': SLATE, 'space_after': 4.8,
                   'line': 1.06} for f_ in fields]
        text(slide, x + 0.28, y + 0.66, cw - 0.5, ch - 0.86, blocks)
    text(slide, MARGIN, FOOTER_Y - 0.5, CONTENT_W, 0.3,
         [{'runs': [
             ('Integrity where it belongs: ', {'size': 11.5, 'bold': True, 'color': INK}),
             ('unique constraints, foreign keys and cascade rules live in the schema, '
              'so no interface path can create a duplicate enrollment or an orphaned '
              'grade.', {'size': 11.5, 'color': SLATE})]}])
    return slide


def slide_features(prs):
    slide = head(prs, 8, 'Key Features', kicker='What the system actually does')
    features = [
        ('Academic Records', ACCENT,
         'CRUD for students, courses and enrollments. Administrators issue lecturer '
         'accounts on Teachers. Enrollment is refused once max_students is reached.'),
        ('Weighted Gradebook', GREEN,
         'Coursework, midterm, exam and retake — each with a weight. The course mark '
         'is their weighted mean. assigned_by records who last saved the component.'),
        ('Attendance Register', AMBER,
         'Present, absent and late — one row per enrollment per day, enforced in the '
         'schema. Mark Class saves a whole roster in a single request.'),
        ('Analytics & Export', VIOLET,
         'Four KPI tiles, three Chart.js views and JSON at /api/dashboard/. CSV '
         'export of students and grades is scoped to the caller\'s role.'),
        ('Role-Based Access', TEAL,
         'Administrators see the institution; teachers see only their courses. '
         'Decorators guard the route, querysets filter again — a guessed URL returns nothing.'),
        ('i18n & Quality', PINK,
         '308 strings, English ⇄ Polish, locale-aware dates. 80 automated tests and '
         'four CI gates (tests, migrations, deploy check, collectstatic) on every push.'),
    ]
    cw = (CONTENT_W - 2 * 0.22) / 3
    ch = 2.24
    for i, (title, color, body) in enumerate(features):
        row, col = divmod(i, 3)
        x = MARGIN + col * (cw + 0.22)
        y = BODY_TOP + row * (ch + 0.22)
        card(slide, x, y, cw, ch)
        rect(slide, x, y, cw, 0.055, fill=color, radius=0.5)
        text(slide, x + 0.28, y + 0.3, cw - 0.56, ch - 0.5, [
            {'text': title, 'size': 14.5, 'bold': True, 'color': INK,
             'font': HEAD_FONT, 'space_after': 7},
            {'text': body, 'size': 11.3, 'color': SLATE, 'line': 1.18},
        ])
    return slide


def slide_ui(prs):
    slide = head(prs, 9, 'Application Screenshots', kicker='The running interface')
    shot_w = 4.96
    mob_h = shot_w / 1.6
    mob_w = mob_h * 430 / 932
    total = 2 * shot_w + mob_w + 2 * 0.34
    x0 = MARGIN + (CONTENT_W - total) / 2
    y0 = BODY_TOP + 0.22

    for i, (name, caption) in enumerate([
        ('dashboard-light.png',
         'Dashboard — KPI tiles and three Chart.js views'),
        ('students-list.png',
         'Student register — search, filters, pagination'),
    ]):
        shot = SHOTS / name
        if shot.exists():
            picture_frame(slide, x0 + i * (shot_w + 0.34), y0, shot_w, shot,
                          caption=caption)

    mobile = SHOTS / 'dashboard-mobile.png'
    if mobile.exists():
        mx = x0 + 2 * (shot_w + 0.34)
        card(slide, mx - 0.055, y0 - 0.055, mob_w + 0.11, mob_h + 0.11)
        slide.shapes.add_picture(str(mobile), Inches(mx), Inches(y0),
                                 Inches(mob_w), Inches(mob_h))
        text(slide, mx - 0.45, y0 + mob_h + 0.14, mob_w + 0.9, 0.44, [
            {'text': 'Responsive', 'size': 10, 'color': MUTED, 'space_after': 2},
            {'text': 'down to 430 px', 'size': 10, 'color': MUTED},
        ], align=PP_ALIGN.CENTER)

    note_bar(slide, FOOTER_Y - 0.86, [
        ('Every table is role-aware: ', {'size': 12, 'bold': True, 'color': WHITE}),
        ('a teacher opening the same page sees only students in their own courses, '
         'and the action buttons disappear.',
         {'size': 12, 'color': RGBColor(0xD6, 0xDF, 0xEC)}),
    ], height=0.56)
    return slide


def slide_deployment(prs):
    slide = head(prs, 10, 'Deployment & Results',
                 kicker='In production, not on a laptop')
    chip_w = 6.15
    chip = rect(slide, MARGIN, BODY_TOP, chip_w, 0.62, fill=INK, radius=0.14)
    chip.click_action.hyperlink.address = LIVE_HREF
    rect(slide, MARGIN + 0.24, BODY_TOP + 0.19, 0.16, 0.16, fill=GREEN,
         shape=MSO_SHAPE.OVAL)
    text(slide, MARGIN + 0.52, BODY_TOP + 0.06, chip_w - 0.7, 0.5,
         [{'runs': [('LIVE   ', {'size': 10, 'bold': True, 'color': GREEN,
                                 'spacing': 1.4}),
                    (LIVE_HREF, {'size': 11.5, 'bold': True, 'color': WHITE,
                                 'url': LIVE_HREF})]}],
         anchor=MSO_ANCHOR.MIDDLE)

    login_x = MARGIN + chip_w + 0.22
    login_w = CONTENT_W - chip_w - 0.22
    card(slide, login_x, BODY_TOP, login_w, 0.62, fill=WHITE)
    text(slide, login_x + 0.18, BODY_TOP + 0.08, login_w - 0.36, 0.48, [
        {'text': 'DEFENCE LOGIN  (teacher)', 'size': 8.5, 'bold': True,
         'color': MUTED, 'spacing': 0.8, 'space_after': 2},
        {'runs': [
            (f'{DEMO_TEACHER}  /  {DEMO_PASSWORD}',
             {'size': 12, 'bold': True, 'color': INK}),
        ]},
    ])

    tiles = [('25', 'students', ACCENT), ('8', 'courses', VIOLET),
             ('87', 'enrollments', TEAL), ('87', 'grades', AMBER),
             ('€0', 'monthly cost', GREEN)]
    tw = (CONTENT_W - 4 * 0.2) / 5
    for i, (value, label, color) in enumerate(tiles):
        stat_tile(slide, MARGIN + i * (tw + 0.2), BODY_TOP + 0.78, tw, 0.86,
                  value, label, color=color, value_size=24, label_size=10.5)

    pw = (CONTENT_W - 0.28) / 2
    py = BODY_TOP + 1.82
    panel(slide, MARGIN, py, pw, 2.72, 'Delivered', [
        ('Live on Render + Neon —', 'HTTPS, HSTS, managed PostgreSQL'),
        ('Six-model schema —', 'constraints and an audit trail in the database'),
        ('Role isolation —', 'admins and teachers see strictly their own data'),
        ('Bilingual interface —', '308 strings, locale-aware dates and numbers'),
        ('80 tests, four CI gates —', 'green on every push, then build.sh deploys'),
    ], accent=GREEN, size=11.6, gap=8)
    panel(slide, MARGIN + pw + 0.28, py, pw, 2.72, 'Future work', [
        ('Student portal —', 'a third role with read-only access to own results'),
        ('Semester transcripts —', 'generated PDF reports per student'),
        ('Email notifications —', 'alerts on new grades and absences'),
        ('Documented REST API —', 'token-authenticated public interface'),
        ('Two-factor login —', 'TOTP for administrator accounts'),
    ], accent=ACCENT, size=11.6, gap=8)
    return slide


def slide_thanks(prs):
    slide = blank(prs)
    paint(slide, INK)
    rect(slide, 9.2, -1.4, 6.4, 6.4, fill=INK_SOFT, radius=0.5)
    track(slide, MARGIN, 2.72, 0.11, 1.0, color=ACCENT)
    text(slide, MARGIN + 0.34, 2.66, 9.4, 1.2, [
        {'text': 'Thank you for your attention', 'size': 40, 'bold': True,
         'color': WHITE, 'font': HEAD_FONT, 'space_after': 10},
        {'text': 'Questions are welcome.', 'size': 15, 'color': FAINT},
    ])
    rect(slide, MARGIN, 4.36, 7.9, 0.012, fill=RGBColor(0x2C, 0x3A, 0x52),
         shape=MSO_SHAPE.RECTANGLE)
    text(slide, MARGIN, 4.62, 6.0, 0.9, [
        {'text': STUDENT, 'size': 16, 'bold': True, 'color': WHITE,
         'space_after': 4},
        {'text': f'Computer Science · Index: {INDEX_NO}', 'size': 11.5,
         'color': MUTED, 'space_after': 4},
        {'text': LIVE_HREF, 'size': 11.5, 'bold': True,
         'color': ACCENT, 'url': LIVE_HREF},
    ])
    add_logo(slide, LOGO_DARK_V, 9.35, 4.55, 3.2)
    footer(slide, dark=True)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    slide_title(prs)
    slide_agenda(prs)
    slide_problem(prs)
    slide_competitors(prs)
    slide_stack(prs)
    slide_architecture(prs)
    slide_database(prs)
    slide_features(prs)
    slide_ui(prs)
    slide_deployment(prs)
    slide_thanks(prs)

    prs.save(OUTPUT)
    print(f'Wrote {OUTPUT.relative_to(ROOT)} — {len(prs.slides._sldIdLst)} slides')


if __name__ == '__main__':
    main()
